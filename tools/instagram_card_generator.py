"""
Instagram 카드뉴스 생성 Tool

입력: JSON 스펙 (제목, 카드 리스트, 메타데이터)
출력: 10장짜리 .pptx 파일 (1080x1350 비율, Instagram 캐러셀용)

규칙:
- 1장: 표지 카드 (이모지 + 제목 + 후킹)
- 2~9장: 본문 카드 (작성자가 JSON으로 제공)
- 10장: 광고 카드 (고정 템플릿 — "서박사의 영양공식" 브랜드)

사용자가 제공하는 본문 카드는 최대 8장까지. 부족하면 남는 슬롯은 건너뜀.
광고 카드는 항상 마지막에 자동으로 추가됨. 수정 불가.

사용 예:
    python tools/instagram_card_generator.py \
        --input .tmp/instagram/omega3_spec.json \
        --output "대시보드/서박사의 영양공식/인스타/2026-04-09_omega3.pptx"

JSON 스펙 형식:
{
    "post_title": "전 세계 76%가 오메가-3 부족 — 한국인은 어떨까?",
    "cover_emoji": "🌊",
    "cover_hook": "전 세계 76%가 오메가-3 부족",
    "cover_subtitle": "한국인은 어떨까?",
    "theme": "ocean",      // 선택: ocean | sage | terracotta | berry | coral
    "cards": [
        {
            "type": "stat",         // stat | list | quote | body
            "eyebrow": "최신 연구",
            "headline": "76%",
            "unit": "%",            // stat 전용: 큰 숫자 옆에 작게 붙음
            "body": "세계 인구가 EPA+DHA 권장 섭취량에 미달"
        },
        {
            "type": "list",
            "eyebrow": "좋은 제품 고르는 법",
            "headline": "체크리스트 4가지",
            "items": [
                "EPA+DHA 합산 함량 확인",
                "산패도 (IFOS 인증)",
                "원료 투명성 (어종·어장)",
                "냉장 보관 가능한 포장"
            ]
        },
        {
            "type": "body",
            "eyebrow": "한국인의 현실",
            "headline": "생선 섭취량, 계속 감소 중",
            "body": "전통적으로 생선을 많이 먹는다고 알려졌지만, 국민건강영양조사에 따르면 최근 수십 년간 한국인의 어류 섭취량은 꾸준히 줄어왔습니다."
        }
    ]
}
"""

import argparse
import json
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn


# ============================================================================
# 디자인 상수
# ============================================================================

# Instagram 캐러셀 포트레이트: 1080 x 1350 (4:5)
# python-pptx는 inch 단위. 10" x 12.5" = 4:5 비율
SLIDE_W = Inches(10)
SLIDE_H = Inches(12.5)

# 테마 색상 팔레트
THEMES = {
    "ocean": {
        "bg_dark":   "#065A82",   # 딥 블루
        "bg_light":  "#F2F8FC",   # 아주 연한 블루
        "accent":    "#1C7293",   # 틸
        "deep":      "#21295C",   # 미드나잇
        "text_on_dark":  "#FFFFFF",
        "text_on_light": "#1A2238",
        "muted":         "#6C7A93",
    },
    "sage": {
        "bg_dark":   "#2C5F2D",
        "bg_light":  "#F5F9F3",
        "accent":    "#97BC62",
        "deep":      "#1A3C1B",
        "text_on_dark":  "#FFFFFF",
        "text_on_light": "#1A3C1B",
        "muted":         "#6B7F5E",
    },
    "terracotta": {
        "bg_dark":   "#B85042",
        "bg_light":  "#FBF6F1",
        "accent":    "#A7BEAE",
        "deep":      "#6B2A20",
        "text_on_dark":  "#FFFFFF",
        "text_on_light": "#3A1A14",
        "muted":         "#8B6F65",
    },
    "berry": {
        "bg_dark":   "#6D2E46",
        "bg_light":  "#FCF7F2",
        "accent":    "#A26769",
        "deep":      "#3D1829",
        "text_on_dark":  "#FFFFFF",
        "text_on_light": "#2E0F1C",
        "muted":         "#9A7C7E",
    },
    "coral": {
        "bg_dark":   "#F96167",
        "bg_light":  "#FFF9EC",
        "accent":    "#2F3C7E",
        "deep":      "#8A1B20",
        "text_on_dark":  "#FFFFFF",
        "text_on_light": "#2F3C7E",
        "muted":         "#7B8097",
    },
}

# 한글 폰트 (Windows/Mac/Linux 모두 존재할 가능성 높은 순서)
FONT_BOLD = "Malgun Gothic"   # 맑은 고딕 (Windows 기본)
FONT_BODY = "Malgun Gothic"

# 브랜드 상수
BRAND_NAME = "서박사의 영양공식"
BRAND_TAGLINE = "3분 안에 찾는 나만의 맞춤 영양제"
BRAND_URL = "nutriformula.co.kr"
BRAND_CTA = "프로필 링크에서 무료 분석 →"


# ============================================================================
# 유틸리티
# ============================================================================

def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_rect(slide, left, top, width, height, fill_hex, line_hex=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    if line_hex is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_to_rgb(line_hex)
        shape.line.width = Pt(0)
    shape.shadow.inherit = False
    return shape


def add_rounded(slide, left, top, width, height, fill_hex, line_hex=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.15
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    if line_hex is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_to_rgb(line_hex)
        shape.line.width = Pt(1)
    return shape


def add_text(slide, left, top, width, height, text,
             font_size=16, bold=False, color_hex="#000000",
             align="left", font_name=FONT_BODY, line_spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    lines = text.split("\n") if isinstance(text, str) else [str(text)]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }.get(align, PP_ALIGN.LEFT)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = hex_to_rgb(color_hex)
    return tb


def add_page_number(slide, index, total, theme, on_dark=False):
    """우상단 페이지 번호 '03 / 10' 형식"""
    txt = f"{index:02d} / {total:02d}"
    color = theme["text_on_dark"] if on_dark else theme["muted"]
    add_text(
        slide,
        left=Inches(8.5), top=Inches(0.5),
        width=Inches(1.3), height=Inches(0.4),
        text=txt, font_size=12, bold=True, color_hex=color,
        align="right", font_name=FONT_BODY,
    )


def add_brand_footer(slide, theme, on_dark=False):
    """좌하단 브랜드 푸터 '@서박사의영양공식'"""
    color = theme["text_on_dark"] if on_dark else theme["muted"]
    add_text(
        slide,
        left=Inches(0.5), top=Inches(11.7),
        width=Inches(6), height=Inches(0.4),
        text=f"@{BRAND_NAME.replace(' ', '')}  |  {BRAND_URL}",
        font_size=11, bold=False, color_hex=color,
        align="left", font_name=FONT_BODY,
    )


# ============================================================================
# 카드 타입별 렌더러
# ============================================================================

def render_cover(slide, spec, theme, total):
    """카드 1: 표지 — 큰 이모지 + 후킹 헤드라인"""
    # 전체 배경: 다크
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, theme["bg_dark"])

    # 상단 얇은 악센트 바
    add_rect(slide, Emu(0), Inches(0.0), SLIDE_W, Inches(0.15), theme["accent"])

    # 상단 eyebrow
    add_text(
        slide, Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.5),
        text="SEOBAKSA NUTRITION  |  WEEKLY",
        font_size=13, bold=True,
        color_hex=theme["accent"], align="left",
    )

    # 이모지 (아주 크게)
    emoji = spec.get("cover_emoji", "📝")
    add_text(
        slide, Inches(0.8), Inches(2.0), Inches(8.4), Inches(2.5),
        text=emoji, font_size=180, bold=False,
        color_hex=theme["text_on_dark"], align="left",
    )

    # 메인 헤드라인 (후킹)
    hook = spec.get("cover_hook", spec.get("post_title", ""))
    add_text(
        slide, Inches(0.8), Inches(5.0), Inches(8.4), Inches(3.5),
        text=hook, font_size=56, bold=True,
        color_hex=theme["text_on_dark"], align="left",
        font_name=FONT_BOLD, line_spacing=1.15,
    )

    # 서브타이틀
    subtitle = spec.get("cover_subtitle", "")
    if subtitle:
        add_text(
            slide, Inches(0.8), Inches(8.8), Inches(8.4), Inches(1.5),
            text=subtitle, font_size=32, bold=False,
            color_hex=theme["accent"], align="left",
            line_spacing=1.2,
        )

    # 좌하단 "SWIPE →" 인디케이터
    add_text(
        slide, Inches(0.8), Inches(10.9), Inches(4), Inches(0.6),
        text="밀어서 자세히 보기  →",
        font_size=16, bold=True,
        color_hex=theme["text_on_dark"], align="left",
    )

    add_page_number(slide, 1, total, theme, on_dark=True)
    add_brand_footer(slide, theme, on_dark=True)


def render_stat(slide, card, index, total, theme):
    """통계 카드 — 거대한 숫자 + 설명"""
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, theme["bg_light"])

    # 좌측 악센트 바
    add_rect(slide, Emu(0), Emu(0), Inches(0.3), SLIDE_H, theme["accent"])

    # Eyebrow
    if card.get("eyebrow"):
        add_text(
            slide, Inches(0.8), Inches(1.3), Inches(8.4), Inches(0.5),
            text=card["eyebrow"].upper(),
            font_size=14, bold=True,
            color_hex=theme["accent"], align="left",
        )

    # 거대 숫자
    headline = card.get("headline", "")
    unit = card.get("unit", "")
    add_text(
        slide, Inches(0.8), Inches(2.5), Inches(8.4), Inches(4.5),
        text=headline, font_size=280, bold=True,
        color_hex=theme["bg_dark"], align="left",
        font_name=FONT_BOLD, line_spacing=0.9,
    )
    if unit and unit not in headline:
        add_text(
            slide, Inches(6), Inches(3), Inches(3), Inches(2),
            text=unit, font_size=80, bold=True,
            color_hex=theme["accent"], align="left",
        )

    # 설명
    body = card.get("body", "")
    if body:
        add_text(
            slide, Inches(0.8), Inches(7.5), Inches(8.4), Inches(3),
            text=body, font_size=28, bold=False,
            color_hex=theme["text_on_light"], align="left",
            line_spacing=1.4,
        )

    add_page_number(slide, index, total, theme, on_dark=False)
    add_brand_footer(slide, theme, on_dark=False)


def render_list(slide, card, index, total, theme):
    """리스트 카드 — 헤드라인 + 항목들"""
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, theme["bg_light"])
    add_rect(slide, Emu(0), Emu(0), Inches(0.3), SLIDE_H, theme["accent"])

    if card.get("eyebrow"):
        add_text(
            slide, Inches(0.8), Inches(1.3), Inches(8.4), Inches(0.5),
            text=card["eyebrow"].upper(),
            font_size=14, bold=True,
            color_hex=theme["accent"], align="left",
        )

    add_text(
        slide, Inches(0.8), Inches(2.0), Inches(8.4), Inches(1.8),
        text=card.get("headline", ""), font_size=44, bold=True,
        color_hex=theme["bg_dark"], align="left",
        font_name=FONT_BOLD, line_spacing=1.15,
    )

    # 항목들: 큰 숫자 배지 + 텍스트
    items = card.get("items", [])[:5]
    top = 5.0
    row_h = 1.2
    for i, item in enumerate(items):
        # 번호 배지 (원)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8), Inches(top + i * row_h),
            Inches(0.9), Inches(0.9),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = hex_to_rgb(theme["bg_dark"])
        circle.line.fill.background()
        circle_tf = circle.text_frame
        circle_tf.margin_left = Emu(0)
        circle_tf.margin_right = Emu(0)
        circle_tf.margin_top = Emu(0)
        circle_tf.margin_bottom = Emu(0)
        circle_p = circle_tf.paragraphs[0]
        circle_p.alignment = PP_ALIGN.CENTER
        circle_run = circle_p.add_run()
        circle_run.text = str(i + 1)
        circle_run.font.name = FONT_BOLD
        circle_run.font.size = Pt(22)
        circle_run.font.bold = True
        circle_run.font.color.rgb = hex_to_rgb(theme["text_on_dark"])

        # 텍스트
        add_text(
            slide, Inches(2.0), Inches(top + i * row_h + 0.15),
            Inches(7.5), Inches(1),
            text=item, font_size=24, bold=True,
            color_hex=theme["text_on_light"], align="left",
            line_spacing=1.2,
        )

    add_page_number(slide, index, total, theme, on_dark=False)
    add_brand_footer(slide, theme, on_dark=False)


def render_body(slide, card, index, total, theme):
    """본문 카드 — 헤드라인 + 긴 본문 텍스트"""
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, theme["bg_light"])
    add_rect(slide, Emu(0), Emu(0), Inches(0.3), SLIDE_H, theme["accent"])

    if card.get("eyebrow"):
        add_text(
            slide, Inches(0.8), Inches(1.3), Inches(8.4), Inches(0.5),
            text=card["eyebrow"].upper(),
            font_size=14, bold=True,
            color_hex=theme["accent"], align="left",
        )

    add_text(
        slide, Inches(0.8), Inches(2.0), Inches(8.4), Inches(2.5),
        text=card.get("headline", ""), font_size=44, bold=True,
        color_hex=theme["bg_dark"], align="left",
        font_name=FONT_BOLD, line_spacing=1.15,
    )

    add_text(
        slide, Inches(0.8), Inches(5.0), Inches(8.4), Inches(6),
        text=card.get("body", ""), font_size=24, bold=False,
        color_hex=theme["text_on_light"], align="left",
        line_spacing=1.55,
    )

    add_page_number(slide, index, total, theme, on_dark=False)
    add_brand_footer(slide, theme, on_dark=False)


def render_quote(slide, card, index, total, theme):
    """인용 카드 — 큰따옴표 + 인용문"""
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, theme["bg_dark"])

    # 큰 따옴표
    add_text(
        slide, Inches(0.8), Inches(1.5), Inches(3), Inches(3),
        text="\u201C", font_size=320, bold=True,
        color_hex=theme["accent"], align="left",
    )

    add_text(
        slide, Inches(0.8), Inches(4.5), Inches(8.4), Inches(5.5),
        text=card.get("body", ""), font_size=40, bold=True,
        color_hex=theme["text_on_dark"], align="left",
        font_name=FONT_BOLD, line_spacing=1.3,
    )

    if card.get("eyebrow"):
        add_text(
            slide, Inches(0.8), Inches(10.3), Inches(8.4), Inches(0.5),
            text=f"— {card['eyebrow']}",
            font_size=18, bold=False,
            color_hex=theme["accent"], align="left",
        )

    add_page_number(slide, index, total, theme, on_dark=True)
    add_brand_footer(slide, theme, on_dark=True)


def render_ad(slide, theme, total, logo_path=None):
    """
    마지막 고정 광고 카드 — "서박사의 영양공식"
    이 카드는 절대 수정되지 않음. 입력 무시.
    """
    # 전체 배경: 딥 컬러
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, theme["deep"])

    # 상단 악센트 바
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.15), theme["accent"])

    # Eyebrow
    add_text(
        slide, Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.5),
        text="MADE BY",
        font_size=14, bold=True,
        color_hex=theme["accent"], align="center",
    )

    # 로고 (있으면 배치)
    if logo_path and os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(
                logo_path,
                Inches(4.0), Inches(2.0),
                width=Inches(2), height=Inches(2),
            )
        except Exception:
            pass

    # 브랜드명
    add_text(
        slide, Inches(0.8), Inches(4.3), Inches(8.4), Inches(1.5),
        text=BRAND_NAME,
        font_size=72, bold=True,
        color_hex=theme["text_on_dark"], align="center",
        font_name=FONT_BOLD,
    )

    # 태그라인
    add_text(
        slide, Inches(0.8), Inches(6.0), Inches(8.4), Inches(1),
        text=BRAND_TAGLINE,
        font_size=28, bold=False,
        color_hex=theme["accent"], align="center",
    )

    # 구분선
    add_rect(
        slide,
        Inches(4.0), Inches(7.4),
        Inches(2), Inches(0.04),
        theme["accent"],
    )

    # 핵심 설명
    add_text(
        slide, Inches(0.8), Inches(7.8), Inches(8.4), Inches(2),
        text="36가지 설문과 라이프스타일을 분석해\n과학적 근거에 기반한 영양제를 추천합니다.",
        font_size=22, bold=False,
        color_hex=theme["text_on_dark"], align="center",
        line_spacing=1.5,
    )

    # CTA 박스
    cta_top = Inches(10.0)
    add_rounded(
        slide,
        Inches(1.5), cta_top,
        Inches(7), Inches(1.2),
        theme["accent"],
    )
    add_text(
        slide, Inches(1.5), cta_top + Inches(0.3),
        Inches(7), Inches(0.7),
        text=BRAND_CTA,
        font_size=26, bold=True,
        color_hex=theme["deep"], align="center",
        font_name=FONT_BOLD,
    )

    # 하단 URL
    add_text(
        slide, Inches(0.8), Inches(11.7), Inches(8.4), Inches(0.5),
        text=BRAND_URL,
        font_size=18, bold=True,
        color_hex=theme["text_on_dark"], align="center",
    )

    # 페이지 번호
    add_text(
        slide, Inches(8.5), Inches(0.5),
        Inches(1.3), Inches(0.4),
        text=f"{total:02d} / {total:02d}",
        font_size=12, bold=True,
        color_hex=theme["accent"], align="right",
    )


# ============================================================================
# 메인 빌더
# ============================================================================

def build_carousel(spec: dict, output_path: str, logo_path: str = None):
    theme_name = spec.get("theme", "ocean")
    theme = THEMES.get(theme_name, THEMES["ocean"])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 총 슬라이드 수 = 표지 + 본문(최대 8) + 광고 = 10
    body_cards = (spec.get("cards") or [])[:8]
    total = 1 + len(body_cards) + 1   # 표지 + 본문 + 광고
    blank_layout = prs.slide_layouts[6]

    # 1. 표지
    cover_slide = prs.slides.add_slide(blank_layout)
    render_cover(cover_slide, spec, theme, total)

    # 2~N-1. 본문
    for i, card in enumerate(body_cards):
        slide = prs.slides.add_slide(blank_layout)
        card_type = card.get("type", "body")
        index = i + 2
        if card_type == "stat":
            render_stat(slide, card, index, total, theme)
        elif card_type == "list":
            render_list(slide, card, index, total, theme)
        elif card_type == "quote":
            render_quote(slide, card, index, total, theme)
        else:
            render_body(slide, card, index, total, theme)

    # N. 광고 (고정)
    ad_slide = prs.slides.add_slide(blank_layout)
    render_ad(ad_slide, theme, total, logo_path=logo_path)

    # 저장
    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out.resolve()), total


def main():
    parser = argparse.ArgumentParser(description="Instagram 카드뉴스 PPTX 생성")
    parser.add_argument("--input", "-i", required=True, help="JSON 스펙 파일 경로")
    parser.add_argument("--output", "-o", required=True, help="출력 PPTX 경로")
    parser.add_argument("--logo", default=None, help="광고 카드에 넣을 로고 이미지 (선택)")
    args = parser.parse_args()

    with open(args.input, "r", encoding="utf-8") as f:
        spec = json.load(f)

    path, total = build_carousel(spec, args.output, logo_path=args.logo)
    print(f"OK  {total} slides → {path}")


if __name__ == "__main__":
    main()
